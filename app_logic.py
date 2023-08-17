import os
import tempfile
import atexit
import fitz
import docx
from pptx import Presentation
from docx import Document
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml

class Logic:
    def __init__(self):
        """
        初始化逻辑类实例。

        使用 Logic 类可以进行文件搜索、内容检索等操作。
        """
        self.found_files = []
        self.temp_folder = os.path.join(tempfile.gettempdir(), "highlight_temp")
        os.makedirs(self.temp_folder, exist_ok=True)

    def search_files(self, folder_path, keywords):
        """
        在给定的文件夹路径中递归搜索包含指定关键词的文件。

        Args:
            folder_path (str): 要搜索的文件夹路径。
            keywords (list): 要搜索的关键词列表。

        Returns:
            list: 包含符合搜索条件的文件路径及关键词列表的元组列表。
        """
        self.found_files = []

        def search_recursively(folder):
            """
            递归搜索文件夹中包含关键词的文件。

            Args:
                folder (str): 要搜索的文件夹路径。
            """
            for root, dirs, files in os.walk(folder):
                for file in files:
                    try:
                        file_path = os.path.join(root, file)
                        norm_file_path = os.path.normpath(file_path)
                        if not file.startswith(('~', '.', '$')):
                            if file.endswith('.txt'):
                                with open(norm_file_path, 'rb') as f:
                                    content = f.read().decode('utf-8', errors='ignore')
                                    key = []
                                    for keyword in keywords:
                                            if keyword in content:
                                                key.append(keyword)
                                    if len(key) > 0:
                                            self.found_files.append((norm_file_path,key))
                            elif file.endswith('.docx'):
                                doc = docx.Document(norm_file_path)
                                content = '\n'.join(para.text for para in doc.paragraphs)
                                key = []
                                for keyword in keywords:
                                        if keyword in content:
                                            key.append(keyword)
                                if len(key) > 0:
                                        self.found_files.append((norm_file_path,key))
                            elif file.endswith('.pptx'):
                                prs = Presentation(norm_file_path)
                                content = ''
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if shape.has_text_frame:
                                            for paragraph in shape.text_frame.paragraphs:
                                                content += paragraph.text + '\n'
                                key = []
                                for keyword in keywords:
                                        if keyword in content:
                                            key.append(keyword)
                                if len(key) > 0:
                                        self.found_files.append((norm_file_path,key))
                            elif file.endswith('.pdf'):
                                content = self.pdf_search(norm_file_path, keywords)
                                key = []
                                for keyword in keywords:
                                        if keyword in content:
                                            key.append(keyword)
                                if len(key) > 0:
                                        self.found_files.append((norm_file_path,key))
                    except Exception as e:
                        print(f"Error processing {file_path}: {e}")

        search_recursively(folder_path)
        return self.found_files

    def pdf_search(self, file_path, keywords):
        """
        在给定的 PDF 文件中搜索包含指定关键词的内容。

        Args:
            file_path (str): 要搜索的 PDF 文件路径。
            keywords (list): 要搜索的关键词列表。

        Returns:
            str: 包含指定关键词的内容文本。
        """
        content = ''
        pdf_document = fitz.open(file_path)
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            text = page.get_text()
            content += text + '\n'
        return content

    def highlight_keywords(self, pdf_file, keywords):
        """
        在 PDF 文件中高亮指定关键词。

        Args:
            pdf_file (str): 要高亮的 PDF 文件路径。
            keywords (list): 要高亮的关键词列表。
        """
        try:
            if isinstance(keywords, str):
                keywords = [keywords]
            pdf_document = fitz.open(pdf_file)
            new_pdf_path = os.path.join(self.temp_folder, "highlighted_pdf.pdf")

            new_pdf_document = fitz.open()  # 创建一个新的PDF文档

            new_pdf_document.insert_pdf(pdf_document)
            page_count = []
            for page_num in range(new_pdf_document.page_count):
                page = new_pdf_document[page_num]
                text = page.get_text()
                contains_keyword = any(keyword in text for keyword in keywords)

                if not contains_keyword:
                    page_count.append(page_num)
            new_pdf_document.delete_pages(page_count)

            for keyword in keywords:
                print(keyword)
                for page_num in range(new_pdf_document.page_count):
                    page = new_pdf_document[page_num]
                    hits = page.search_for(keyword)
                    for hit in hits:
                        highlight = page.add_highlight_annot(hit)
                        highlight.update(opacity=0.5)  # 设置高亮的透明度

            new_pdf_document.save(new_pdf_path)  # 保存修改后的PDF
            new_pdf_document.close()
            pdf_document.close()
            os.system(new_pdf_path)
        except Exception as e:
            print(f"Error highlighting keywords: {e}")

    def search_keywords_in_file(self, file_path, keywords):
        """
        在指定文件中搜索包含关键词的内容。

        Args:
            file_path (str): 要搜索的文件路径。
            keywords (list): 要搜索的关键词列表。

        Returns:
            list: 包含关键词的内容行列表的元组。
        """
        try:
            if isinstance(keywords, str):
                keywords = [keywords]

            if file_path.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif file_path.endswith('.docx'):
                doc = docx.Document(file_path)
                content = '\n'.join(para.text for para in doc.paragraphs)
            elif file_path.endswith('.pptx'):
                prs = Presentation(file_path)
                content = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                content += paragraph.text + '\n'
            elif file_path.endswith('.pdf'):
                content = self.pdf_search(file_path, keywords)
            else:
                content = ''

            content_lines = content.split('\n')
            lines_with_keywords = []
            for line_num, line in enumerate(content_lines, start=1):
                for keyword in keywords:
                    if keyword in line:
                        lines_with_keywords.append((line_num, line))
                        break

            return lines_with_keywords
        except Exception as e:
            print(f"Error opening {file_path}: {e}")
            return []

    def cleanup_temp_folder(self):
        """
        清理临时文件夹中的临时文件。
        """
        temp_folder = os.path.join(tempfile.gettempdir(), "highlight_temp")
        if os.path.exists(temp_folder):
            for filename in os.listdir(temp_folder):
                file_path = os.path.join(temp_folder, filename)
                os.remove(file_path)
            os.rmdir(temp_folder)

        atexit.register(self.cleanup_temp_folder)
