import os
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
import docx  # pip install python-docx
from pptx import Presentation  # pip install python-pptx
import tempfile
import atexit
import fitz

def search_files(folder_path, keyword):
    """
    在给定的文件夹路径中递归搜索包含指定关键词的文件。

    Args:
        folder_path (str): 要搜索的文件夹路径。
        keyword (str): 要搜索的关键词。

    Returns:
        list: 包含指定关键词的文件路径列表。
    """
    found_files = []

    def search_recursively(folder):
        for root, dirs, files in os.walk(folder):
            for file in files:
                try:
                    file_path = os.path.join(root, file)
                    norm_file_path = os.path.normpath(file_path)
                    if not file.startswith(('~', '.', '$')):
                        if file.endswith('.txt'):
                            with open(norm_file_path, 'rb') as f:
                                content = f.read().decode('utf-8', errors='ignore')
                                if keyword in content:
                                    found_files.append(norm_file_path)
                        elif file.endswith('.docx'):
                            doc = docx.Document(norm_file_path)
                            content = '\n'.join(para.text for para in doc.paragraphs)
                            if keyword in content:
                                found_files.append(norm_file_path)
                        elif file.endswith('.pptx'):
                            prs = Presentation(norm_file_path)
                            content = ''
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if shape.has_text_frame:
                                        for paragraph in shape.text_frame.paragraphs:
                                            content += paragraph.text + '\n'
                            if keyword in content:
                                found_files.append(norm_file_path)
                        elif file.endswith('.pdf'):
                            content = pdf_search(norm_file_path, keyword)
                            if keyword in content:
                                found_files.append(norm_file_path)
                except Exception as e:
                    print(f"Error processing {file_path}: {e}")

    search_recursively(folder_path)
    return found_files

def pdf_search(file_path, keywords):
    """
    在给定的 PDF 文件中搜索包含指定关键词的内容。

    Args:
        file_path (str): 要搜索的 PDF 文件路径。
        keywords (str): 要搜索的关键词。

    Returns:
        str: 包含指定关键词的内容文本。
    """
    content = ''
    pdf_document = fitz.open(file_path)
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        text = page.get_text()
        for keyword in keywords:
            if keyword in text:
                content += text + '\n'
    return content

def highlight_keywords(pdf_file, keywords):
    """
    在 PDF 文件中高亮指定关键词。

    Args:
        pdf_file (str): 要高亮的 PDF 文件路径。
        keywords (list): 要高亮的关键词列表。
    """
    try:
        pdf_document = fitz.open(pdf_file)
        
        temp_folder = os.path.join(tempfile.gettempdir(), "pdf_highlight_temp")
        os.makedirs(temp_folder, exist_ok=True)
        new_pdf_path = os.path.join(temp_folder, "highlighted_pdf.pdf")

        new_pdf_document = fitz.open()  # 创建一个新的 PDF 文档

        new_pdf_document.insert_pdf(pdf_document)
        page_count = []
        for page_num in range(new_pdf_document.page_count):
            page = new_pdf_document[page_num]
            text = page.get_text()
            contains_keyword = any(keyword in text for keyword in keywords)
            
            if not contains_keyword:
                page_count.append(page_num)
        new_pdf_document.delete_pages(page_count)

        for keyword in keywords:
            for page_num in range(new_pdf_document.page_count):
                page = new_pdf_document[page_num]
                hits = page.search_for(keyword)
                for hit in hits:
                    highlight = page.add_highlight_annot(hit)
                    highlight.update(opacity=0.5)  # 设置高亮的透明度

        new_pdf_document.save(new_pdf_path)  # 保存修改后的 PDF
        new_pdf_document.close()
        pdf_document.close()
        os.system(new_pdf_path)
    except Exception as e:
        print(f"Error highlighting keywords: {e}")

def show_content():
    """
    在结果文本框中显示选定文件的内容或高亮关键词。
    """
    selected_item = result_listbox.curselection()
    if not selected_item:
        return

    selected_index = selected_item[0]
    selected_file = found_files[selected_index]
    keywords = keyword_entry.get().split(',')
    if open_with_highlight_var.get() and selected_file.endswith('.pdf'):
        highlight_keywords(selected_file, keywords)
    else:
        lines_with_keywords = open_and_search(selected_file, keywords)
        result_text.delete(1.0, tk.END)  # 清空结果文本框
        for line_num, line in lines_with_keywords:
            result_text.insert(tk.END, f"{line_num}: {line}\n")

def open_and_search(file_path, keywords):
    """
    打开指定文件并搜索包含关键词的内容。

    Args:
        file_path (str): 要打开的文件路径。
        keywords (list): 要搜索的关键词列表。

    Returns:
        list: 包含关键词的内容行列表。
    """
    try:
        if isinstance(keywords, str):
            keywords = [keywords]

        if file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif file_path.endswith('.docx'):
            doc = docx.Document(file_path)
            content = '\n'.join(para.text for para in doc.paragraphs)
        elif file_path.endswith('.pptx'):
            prs = Presentation(file_path)
            content = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            content += paragraph.text + '\n'
        elif file_path.endswith('.pdf'):
            content = pdf_search(file_path, keywords)
        else:
            content = ''

        content_lines = content.split('\n')
        lines_with_keywords = []
        for line_num, line in enumerate(content_lines, start=1):
            for keyword in keywords:
                if keyword in line:
                    lines_with_keywords.append((line_num, line))
                    break

        return lines_with_keywords
    except Exception as e:
        print(f"Error opening {file_path}: {e}")
        return []

def search_and_open():
    """
    在指定文件夹中搜索并打开包含关键词的文件。
    """
    folder_path = filedialog.askdirectory(title="选择文件夹")
    keyword = keyword_entry.get()

    if not folder_path or not keyword:
        messagebox.showerror("错误", "请选择文件夹并输入关键词")
        return

    global found_files
    found_files = search_files(folder_path, keyword)

    if not found_files:
        messagebox.showinfo("结果", "未找到包含关键词的文件")
        return

    result_listbox.delete(0, tk.END)
    result_text.delete(1.0, tk.END)

    for file in found_files:
        result_listbox.insert(tk.END, f"在文件 {file} 中找到关键词 '{keyword}'")

def cleanup_temp_folder():
    """
    清理临时文件夹中的临时文件。
    """
    temp_folder = os.path.join(tempfile.gettempdir(), "pdf_highlight_temp")
    if os.path.exists(temp_folder):
        for filename in os.listdir(temp_folder):
            file_path = os.path.join(temp_folder, filename)
            os.remove(file_path)
        os.rmdir(temp_folder)

# 注册清理函数
atexit.register(cleanup_temp_folder)

# 创建主窗口
root = tk.Tk()
root.title("文件搜索和查看工具")

# 创建界面元素
# ... (UI elements remain unchanged)
# 创建界面元素
current_path_label = tk.Label(root, text="当前路径：", anchor='w', padx=10) 
folder_button = tk.Button(root, text="选择文件夹", command=search_and_open)
keyword_label = tk.Label(root, text="请输入关键词（多个关键词用逗号分隔）：")
keyword_entry = tk.Entry(root, width=80)
result_listbox = tk.Listbox(root, selectmode=tk.SINGLE, width=80)
result_text = scrolledtext.ScrolledText(root, wrap=tk.WORD, height=20, width=80)
show_button = tk.Button(root, text="显示内容", command=show_content)

# 布局界面元素
folder_button.pack(pady=10)
keyword_label.pack()
keyword_entry.pack()
result_listbox.pack(padx=10, pady=5)
result_text.pack(padx=10, pady=5)
show_button.pack(pady=5)
current_path_label.pack(fill=tk.X)

# 创建界面元素
# ... (UI elements remain unchanged)
# 创建界面元素
open_with_highlight_var = tk.BooleanVar()  # 创建一个Boolean变量用于控制功能开启
open_with_highlight_checkbox = tk.Checkbutton(root, text="打开文件并高亮关键词", variable=open_with_highlight_var)
open_with_highlight_checkbox.pack(pady=5)

# 运行主循环
root.mainloop()
